import os
from flask import Flask, render_template, request, jsonify, session, Response, stream_with_context, send_file, redirect, url_for, flash
from together import Together # We'll still use this for chat
import requests # For the dedicated file upload
from dotenv import load_dotenv
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
import logging
import io
import csv
import json
import time
import subprocess
import tempfile
import shutil
import sys
import traceback
from pathlib import Path
import uuid
import threading
import signal
from functools import wraps
from collections import defaultdict, deque
import datetime
import base64

# Security and production imports
try:
    from flask_wtf.csrf import CSRFProtect
    from flask_talisman import Talisman
    from flask_limiter import Limiter
    from flask_limiter.util import get_remote_address
    SECURITY_AVAILABLE = True
except ImportError:
    SECURITY_AVAILABLE = False

# Database imports
try:
    from flask_sqlalchemy import SQLAlchemy
    from sqlalchemy import func, text
    from sqlalchemy.orm import relationship
    SQLALCHEMY_AVAILABLE = True
except ImportError:
    SQLALCHEMY_AVAILABLE = False

# Load environment variables from .env file
load_dotenv()

# Rate Limiter Class
class RateLimiter:
    def __init__(self, max_requests=6, time_window=60, max_tokens=60000):
        self.max_requests = max_requests
        self.time_window = time_window
        self.max_tokens = max_tokens
        self.requests = deque()
        self.tokens_used = deque()
        self.lock = threading.Lock()
    
    def can_make_request(self, estimated_tokens=1000):
        with self.lock:
            now = time.time()
            
            # Remove old requests outside time window
            while self.requests and now - self.requests[0] > self.time_window:
                self.requests.popleft()
            
            # Remove old token usage outside time window
            while self.tokens_used and now - self.tokens_used[0][0] > self.time_window:
                self.tokens_used.popleft()
            
            # Check request limit
            if len(self.requests) >= self.max_requests:
                return False, f"Rate limit: {self.max_requests} requests per {self.time_window}s"
            
            # Check token limit
            current_tokens = sum(tokens for _, tokens in self.tokens_used)
            if current_tokens + estimated_tokens > self.max_tokens:
                return False, f"Token limit: {self.max_tokens} tokens per {self.time_window}s"
            
            return True, "OK"
    
    def record_request(self, tokens_used=1000):
        with self.lock:
            now = time.time()
            self.requests.append(now)
            self.tokens_used.append((now, tokens_used))
    
    def get_wait_time(self):
        with self.lock:
            if not self.requests:
                return 0
            return max(0, self.time_window - (time.time() - self.requests[0]))

# Initialize Flask app with production configuration
app = Flask(__name__)

# Load configuration based on environment
from config import get_config
config_class = get_config()
app.config.from_object(config_class)

# Configure logging based on environment
log_level = getattr(logging, app.config['LOG_LEVEL'].upper(), logging.INFO)
logging.basicConfig(
    level=log_level,
    format=app.config['LOG_FORMAT'],
    handlers=[
        logging.FileHandler('businessastra.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Log configuration info
logger.info(f"Starting BusinessAstra in {os.environ.get('FLASK_ENV', 'development')} mode")
logger.info(f"Database: {app.config['SQLALCHEMY_DATABASE_URI'][:50]}...")  # Don't log full DB URL

 
# Initialize database
db = None
if SQLALCHEMY_AVAILABLE:
    db = SQLAlchemy(app)
    logger.info("Database initialized successfully")
else:
    logger.warning("SQLAlchemy not available. Using fallback JSON storage.")

# Initialize security extensions
csrf = None
talisman = None
limiter = None

if SECURITY_AVAILABLE:
    # CSRF Protection
    if app.config.get('WTF_CSRF_ENABLED', False):
        csrf = CSRFProtect(app)
        logger.info("CSRF protection enabled")
    
    # Security headers with Talisman
    if not app.config.get('DEBUG', False):
        csp = {
            'default-src': "'self'",
            'script-src': [
                "'self'",
                "'unsafe-inline'",  # Needed for inline scripts
                'https://cdn.jsdelivr.net',
                'https://cdnjs.cloudflare.com'
            ],
            'style-src': [
                "'self'",
                "'unsafe-inline'",  # Needed for inline styles
                'https://cdn.jsdelivr.net',
                'https://cdnjs.cloudflare.com'
            ],
            'img-src': [
                "'self'",
                'data:',  # For base64 images
                'https:'
            ],
            'font-src': [
                "'self'",
                'https://cdn.jsdelivr.net',
                'https://cdnjs.cloudflare.com'
            ]
        }
        
        talisman = Talisman(
            app,
            force_https=app.config.get('FORCE_HTTPS', False),
            strict_transport_security=True,
            content_security_policy=csp
        )
        logger.info("Security headers configured")
    
    # Rate limiting
    limiter = Limiter(
        app,
        key_func=get_remote_address,
        default_limits=[app.config.get('RATELIMIT_DEFAULT', '100 per hour')],
        storage_uri=app.config.get('RATELIMIT_STORAGE_URL', 'memory://')
    )
    logger.info("Rate limiting configured")
else:
    logger.warning("Security extensions not available. Install flask-wtf, flask-talisman, flask-limiter for production.")

# Document processing libraries
try:
    from docx import Document  # python-docx for Word documents
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available. Word document support disabled.")

try:
    import openpyxl  # For Excel files
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("openpyxl not available. Excel document support disabled.")

try:
    from pptx import Presentation  # python-pptx for PowerPoint
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PowerPoint document support disabled.")

try:
    import pdfplumber  # For PDF text extraction
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logger.warning("pdfplumber not available. PDF text extraction support disabled.")

try:
    from pdf2image import convert_from_bytes  # For PDF to image conversion
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("pdf2image not available. PDF to image conversion support disabled.")

try:
    from PIL import Image  # For image processing
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logger.warning("PIL not available. Image processing support disabled.")

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logger.warning("ReportLab not available. PDF generation support disabled.")

# Database Models
if SQLALCHEMY_AVAILABLE and db:
    class User(db.Model):
        __tablename__ = 'users'
        
        id = db.Column(db.Integer, primary_key=True)
        email = db.Column(db.String(255), unique=True, nullable=False, index=True)
        username = db.Column(db.String(100), nullable=False)
        password_hash = db.Column(db.String(255), nullable=False)
        created_at = db.Column(db.DateTime, default=func.now())
        updated_at = db.Column(db.DateTime, default=func.now(), onupdate=func.now())
        last_login = db.Column(db.DateTime)
        is_active = db.Column(db.Boolean, default=True)
        
        # Relationships
        conversations = relationship("Conversation", back_populates="user", cascade="all, delete-orphan")
        uploads = relationship("FileUpload", back_populates="user", cascade="all, delete-orphan")
        
        def __repr__(self):
            return f'<User {self.username}>'
    
    class Category(db.Model):
        __tablename__ = 'categories'
        
        id = db.Column(db.Integer, primary_key=True)
        name = db.Column(db.String(100), nullable=False, unique=True)
        description = db.Column(db.Text)
        color = db.Column(db.String(7), default='#3498db')  # Hex color
        icon = db.Column(db.String(50), default='💬')
        created_at = db.Column(db.DateTime, default=func.now())
        
        # Relationships
        conversations = relationship("Conversation", back_populates="category")
        
        def __repr__(self):
            return f'<Category {self.name}>'
    
    class Conversation(db.Model):
        __tablename__ = 'conversations'
        
        id = db.Column(db.String(36), primary_key=True)  # UUID
        user_id = db.Column(db.Integer, db.ForeignKey('users.id'), nullable=False)
        category_id = db.Column(db.Integer, db.ForeignKey('categories.id'))
        
        title = db.Column(db.String(200), nullable=False)
        summary = db.Column(db.Text)
        created_at = db.Column(db.DateTime, default=func.now())
        updated_at = db.Column(db.DateTime, default=func.now(), onupdate=func.now())
        last_message_at = db.Column(db.DateTime, default=func.now())
        
        # Conversation metadata
        message_count = db.Column(db.Integer, default=0)
        has_code = db.Column(db.Boolean, default=False)
        has_files = db.Column(db.Boolean, default=False)
        has_analysis = db.Column(db.Boolean, default=False)
        
        # Auto-generated tags
        tags = db.Column(db.Text)  # JSON array of tags
        
        # Relationships
        user = relationship("User", back_populates="conversations")
        category = relationship("Category", back_populates="conversations")
        messages = relationship("Message", back_populates="conversation", cascade="all, delete-orphan")
        uploads = relationship("FileUpload", back_populates="conversation")
        
        def __repr__(self):
            return f'<Conversation {self.title}>'
    
    class Message(db.Model):
        __tablename__ = 'messages'
        
        id = db.Column(db.Integer, primary_key=True)
        conversation_id = db.Column(db.String(36), db.ForeignKey('conversations.id'), nullable=False)
        
        role = db.Column(db.String(20), nullable=False)  # 'user', 'assistant', 'system'
        content = db.Column(db.Text, nullable=False)
        created_at = db.Column(db.DateTime, default=func.now())
        
        # Message metadata
        has_code = db.Column(db.Boolean, default=False)
        code_type = db.Column(db.String(50))  # 'python', 'reportlab', 'manim'
        execution_result = db.Column(db.Text)  # JSON of execution results
        
        # Relationships
        conversation = relationship("Conversation", back_populates="messages")
        
        def __repr__(self):
            return f'<Message {self.role}: {self.content[:50]}...>'
    
    class FileUpload(db.Model):
        __tablename__ = 'file_uploads'
        
        id = db.Column(db.Integer, primary_key=True)
        user_id = db.Column(db.Integer, db.ForeignKey('users.id'), nullable=False)
        conversation_id = db.Column(db.String(36), db.ForeignKey('conversations.id'))
        
        original_filename = db.Column(db.String(255), nullable=False)
        file_type = db.Column(db.String(50), nullable=False)
        file_size = db.Column(db.Integer)
        processed_content = db.Column(db.Text)  # JSONL content
        upload_time = db.Column(db.DateTime, default=func.now())
        
        # Processing metadata
        processing_status = db.Column(db.String(50), default='completed')  # 'processing', 'completed', 'failed'
        processing_method = db.Column(db.String(100))  # 'text_extraction', 'image_conversion', etc.
        
        # Relationships
        user = relationship("User", back_populates="uploads")
        conversation = relationship("Conversation", back_populates="uploads")
        
        def __repr__(self):
            return f'<FileUpload {self.original_filename}>'
        
    # Create default categories
    def create_default_categories():
        """Create default conversation categories"""
        default_categories = [
            {'name': 'Data Analysis', 'description': 'Data processing, visualization, and analytics', 'color': '#e74c3c', 'icon': '📊'},
            {'name': 'Business Reports', 'description': 'Professional reports and documentation', 'color': '#3498db', 'icon': '📋'},
            {'name': 'Code Development', 'description': 'Programming, scripts, and automation', 'color': '#2ecc71', 'icon': '💻'},
            {'name': 'Document Processing', 'description': 'File uploads and document analysis', 'color': '#f39c12', 'icon': '📄'},
            {'name': 'Help & Questions', 'description': 'General assistance and inquiries', 'color': '#9b59b6', 'icon': '❓'},
            {'name': 'Troubleshooting', 'description': 'Error fixes and problem solving', 'color': '#e67e22', 'icon': '🔧'},
            {'name': 'Presentations', 'description': 'Charts, graphs, and visual content', 'color': '#1abc9c', 'icon': '📈'},
        ]
        
        for cat_data in default_categories:
            if not Category.query.filter_by(name=cat_data['name']).first():
                category = Category(**cat_data)
                db.session.add(category)
        
        try:
            db.session.commit()
            logger.info("Default categories created successfully")
        except Exception as e:
            db.session.rollback()
            logger.error(f"Error creating default categories: {e}")
    
    def init_database():
        """Initialize database tables and default data"""
        try:
            # Create all tables
            db.create_all()
            logger.info("Database tables created successfully")
            
            # Create default categories
            create_default_categories()
            
            # Migrate existing JSON data if available
            migrate_json_to_sql()
            
        except Exception as e:
            logger.error(f"Error initializing database: {e}")
    
    def migrate_json_to_sql():
        """Migrate existing JSON user data to SQL database"""
        try:
            if not os.path.exists(USERS_FILE):
                logger.info("No JSON file to migrate")
                return
            
            users_data = load_users()
            if not users_data:
                logger.info("No user data to migrate")
                return
            
            migrated_count = 0
            for email, user_data in users_data.items():
                # Check if user already exists in SQL
                existing_user = User.query.filter_by(email=email).first()
                if existing_user:
                    continue
                
                # Create user in SQL database
                new_user = User(
                    email=email,
                    username=user_data.get('username', 'User'),
                    password_hash=user_data.get('password_hash', ''),
                    created_at=datetime.datetime.fromtimestamp(user_data.get('created_at', time.time()))
                )
                
                db.session.add(new_user)
                db.session.flush()  # Get the user ID
                
                # Migrate chat history
                chat_history = user_data.get('chat_history', [])
                for chat in chat_history:
                    try:
                        # Analyze conversation for categorization
                        messages = chat.get('messages', [])
                        category_name, analysis = categorize_conversation(messages)
                        category = Category.query.filter_by(name=category_name).first()
                        
                        # Create conversation
                        conversation = Conversation(
                            id=chat.get('id', str(uuid.uuid4())),
                            user_id=new_user.id,
                            category_id=category.id if category else None,
                            title=chat.get('name', 'Migrated Chat'),
                            summary=f"Migrated from JSON: {len(messages)} messages",
                            created_at=datetime.datetime.fromtimestamp(chat.get('created_at', time.time())),
                            updated_at=datetime.datetime.fromtimestamp(chat.get('updated_at', time.time())),
                            last_message_at=datetime.datetime.fromtimestamp(chat.get('updated_at', time.time())),
                            message_count=len(messages),
                            has_code=analysis.get('has_code', False),
                            has_files=analysis.get('has_files', False),
                            has_analysis=analysis.get('has_analysis', False),
                            tags=json.dumps(analysis.get('code_types', []))
                        )
                        
                        db.session.add(conversation)
                        db.session.flush()  # Get the conversation ID
                        
                        # Create messages
                        for msg in messages:
                            message = Message(
                                conversation_id=conversation.id,
                                role=msg.get('role', 'user'),
                                content=msg.get('content', ''),
                                has_code='```' in msg.get('content', ''),
                                code_type='python' if '```python' in msg.get('content', '') else None
                            )
                            db.session.add(message)
                    
                    except Exception as e:
                        logger.warning(f"Error migrating chat {chat.get('id', 'unknown')}: {e}")
                        continue
                
                migrated_count += 1
            
            db.session.commit()
            logger.info(f"Successfully migrated {migrated_count} users from JSON to SQL")
            
            # Backup the JSON file
            backup_file = f"{USERS_FILE}.backup.{int(time.time())}"
            shutil.copy2(USERS_FILE, backup_file)
            logger.info(f"JSON file backed up as {backup_file}")
            
        except Exception as e:
            db.session.rollback()
            logger.error(f"Error migrating JSON data: {e}")

else:
    logger.warning("Database models not available. Using fallback storage.")
    
    def init_database():
        """Fallback - no database initialization needed"""
        logger.info("Using JSON fallback storage - no database initialization")
    
    def migrate_json_to_sql():
        """Fallback - no migration possible"""
        pass

# Initialize rate limiter
rate_limiter = RateLimiter(max_requests=5, time_window=60, max_tokens=50000)  # Conservative limits

# Create temporary uploads folder
UPLOAD_FOLDER_PATH = os.path.abspath(app.config['UPLOAD_FOLDER'])
if not os.path.exists(UPLOAD_FOLDER_PATH):
    os.makedirs(UPLOAD_FOLDER_PATH)
    logger.info(f"Created temporary upload folder at {UPLOAD_FOLDER_PATH}")

# --- Together AI Client Initialization ---
TOGETHER_API_KEY = os.getenv("TOGETHER_API_KEY")
if TOGETHER_API_KEY:
    # Clean the API key of any whitespace/newlines
    TOGETHER_API_KEY = TOGETHER_API_KEY.strip().replace('\n', '').replace('\r', '')
    logger.debug(f"API key loaded, length: {len(TOGETHER_API_KEY)}")

TOGETHER_CHAT_CLIENT = None
if TOGETHER_API_KEY:
    try:
        TOGETHER_CHAT_CLIENT = Together(api_key=TOGETHER_API_KEY)
        logger.info("Together AI CHAT client initialized successfully.")
    except Exception as e:
        logger.error(f"Failed to initialize Together AI CHAT client: {e}")
else:
    logger.warning("TOGETHER_API_KEY not found or invalid. AI features will be disabled.")

# --- User Management ---
USERS_FILE = 'users.json'  # Kept for fallback compatibility

def load_users():
    """Load users from JSON file (fallback only)"""
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, 'r') as f:
                return json.load(f)
        except (json.JSONDecodeError, FileNotFoundError):
            return {}
    return {}

def save_users(users):
    """Save users to JSON file (fallback only)"""
    with open(USERS_FILE, 'w') as f:
        json.dump(users, f, indent=2)

def create_user(email, password, username):
    """Create a new user"""
    if SQLALCHEMY_AVAILABLE and db:
        # SQL-based user creation
        try:
            # Check if user already exists
            existing_user = User.query.filter_by(email=email).first()
            if existing_user:
                return False, "User already exists"
            
            # Check if username is already taken
            existing_username = User.query.filter_by(username=username).first()
            if existing_username:
                return False, "Username already taken"
            
            # Create new user
            new_user = User(
                email=email,
                username=username,
                password_hash=generate_password_hash(password)
            )
            
            db.session.add(new_user)
            db.session.commit()
            logger.info(f"Created new user: {username} ({email})")
            return True, "User created successfully"
            
        except Exception as e:
            db.session.rollback()
            logger.error(f"Error creating user: {e}")
            return False, f"Error creating user: {str(e)}"
    else:
        # Fallback to JSON
        users = load_users()
        if email in users:
            return False, "User already exists"
        
        # Check if username is already taken
        for user_data in users.values():
            if user_data.get('username', '').lower() == username.lower():
                return False, "Username already taken"
        
        users[email] = {
            'username': username,
            'password_hash': generate_password_hash(password),
            'created_at': time.time(),
            'chat_history': []  # Initialize empty chat history
        }
        save_users(users)
        return True, "User created successfully"

def verify_user(email, password):
    """Verify user credentials"""
    if SQLALCHEMY_AVAILABLE and db:
        # SQL-based verification
        try:
            user = User.query.filter_by(email=email).first()
            if not user:
                return False
            
            is_valid = check_password_hash(user.password_hash, password)
            if is_valid:
                # Update last login
                user.last_login = func.now()
                db.session.commit()
                logger.info(f"User logged in: {user.username}")
            
            return is_valid
        except Exception as e:
            logger.error(f"Error verifying user: {e}")
            return False
    else:
        # Fallback to JSON
        users = load_users()
        if email not in users:
            return False
        
        return check_password_hash(users[email]['password_hash'], password)

def get_user_by_email(email):
    """Get user object by email"""
    if SQLALCHEMY_AVAILABLE and db:
        return User.query.filter_by(email=email).first()
    else:
        # Fallback to JSON
        users = load_users()
        if email in users:
            user_data = users[email]
            return {
                'email': email,
                'username': user_data.get('username', 'User'),
                'created_at': user_data.get('created_at'),
                'chat_history': user_data.get('chat_history', [])
            }
        return None

def login_required(f):
    """Decorator to require login for routes"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_email' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

# --- Enhanced Conversation Management ---
def get_user_chat_history(user_email):
    """Get chat history for a user with enhanced categorization"""
    if SQLALCHEMY_AVAILABLE and db:
        try:
            user = User.query.filter_by(email=user_email).first()
            if not user:
                return []
            
            # Get conversations with categories, ordered by last message
            conversations = Conversation.query.filter_by(user_id=user.id)\
                .order_by(Conversation.last_message_at.desc())\
                .limit(50).all()
            
            result = []
            for conv in conversations:
                result.append({
                    'id': conv.id,
                    'name': conv.title,
                    'category': {
                        'name': conv.category.name if conv.category else 'General',
                        'color': conv.category.color if conv.category else '#3498db',
                        'icon': conv.category.icon if conv.category else '💬'
                    },
                    'summary': conv.summary,
                    'created_at': conv.created_at.timestamp(),
                    'updated_at': conv.updated_at.timestamp(),
                    'last_message_at': conv.last_message_at.timestamp(),
                    'message_count': conv.message_count,
                    'has_code': conv.has_code,
                    'has_files': conv.has_files,
                    'has_analysis': conv.has_analysis,
                    'tags': json.loads(conv.tags) if conv.tags else []
                })
            
            return result
        except Exception as e:
            logger.error(f"Error getting chat history: {e}")
            return []
    else:
        # Fallback to JSON
        users = load_users()
        user_data = users.get(user_email, {})
        return user_data.get('chat_history', [])

def save_chat_to_history(user_email, chat_data):
    """Save a chat session to user's history"""
    try:
        users = load_users()
        if user_email not in users:
            logger.error(f"User {user_email} not found in database")
            return False
        
        if 'chat_history' not in users[user_email]:
            users[user_email]['chat_history'] = []
        
        existing_chats = users[user_email]['chat_history']
        chat_id = chat_data.get('id')
        
        # Check if this chat ID already exists (prevent duplicates)
        for i, existing_chat in enumerate(existing_chats):
            if existing_chat.get('id') == chat_id:
                # Update existing chat instead of creating duplicate
                existing_chats[i] = chat_data
                save_users(users)
                logger.debug(f"Updated existing chat: {chat_data.get('name')}")
                return True
        
        # Check for similar content to prevent near-duplicates
        chat_messages = chat_data.get('messages', [])
        user_messages = [msg.get('content', '') for msg in chat_messages if msg.get('role') == 'user']
        
        if len(user_messages) > 0:
            first_user_message = user_messages[0][:100]  # First 100 chars
            
            for existing_chat in existing_chats[-5:]:  # Check last 5 chats only
                existing_messages = existing_chat.get('messages', [])
                existing_user_messages = [msg.get('content', '') for msg in existing_messages if msg.get('role') == 'user']
                
                if (len(existing_user_messages) > 0 and 
                    existing_user_messages[0][:100] == first_user_message):
                    logger.debug(f"Skipping similar chat: {chat_data.get('name')}")
                    return False
        
        # Limit chat history to last 30 chats for performance
        if len(existing_chats) >= 30:
            users[user_email]['chat_history'] = existing_chats[-29:]  # Keep last 29, add new one
        
        # Add new chat
        users[user_email]['chat_history'].append(chat_data)
        save_users(users)
        logger.info(f"Saved new chat to history: {chat_data.get('name')} (Total: {len(users[user_email]['chat_history'])})")
        return True
        
    except Exception as e:
        logger.error(f"Error saving chat to history: {e}")
        return False

def categorize_conversation(messages):
    """Enhanced conversation categorization with smart analysis"""
    try:
        # Analyze all messages for better categorization
        content_analysis = {
            'has_code': False,
            'has_files': False,
            'has_analysis': False,
            'keywords': [],
            'code_types': []
        }
        
        full_content = ""
        for msg in messages:
            if msg.get('role') in ['user', 'assistant']:
                content = msg.get('content', '').lower()
                full_content += content + " "
                
                # Check for code
                if any(marker in content for marker in ['```python', '```reportlab', '```manim', 'def ', 'import ', 'class ']):
                    content_analysis['has_code'] = True
                    if '```python' in content or 'def ' in content:
                        content_analysis['code_types'].append('python')
                    if '```reportlab' in content:
                        content_analysis['code_types'].append('reportlab')
                    if '```manim' in content:
                        content_analysis['code_types'].append('manim')
                
                # Check for file uploads
                if any(word in content for word in ['upload', 'file', 'document', 'csv', 'excel', 'pdf']):
                    content_analysis['has_files'] = True
                
                # Check for analysis keywords
                if any(word in content for word in ['analyze', 'analysis', 'data', 'chart', 'graph', 'visualization']):
                    content_analysis['has_analysis'] = True
        
        # Determine category based on content analysis
        category_name = "Help & Questions"  # Default
        
        if content_analysis['has_analysis'] and 'chart' in full_content or 'graph' in full_content:
            category_name = "Presentations"
        elif content_analysis['has_analysis'] and any(word in full_content for word in ['data', 'csv', 'excel']):
            category_name = "Data Analysis"
        elif content_analysis['has_code'] and 'reportlab' in content_analysis['code_types']:
            category_name = "Business Reports"
        elif content_analysis['has_code']:
            category_name = "Code Development"
        elif content_analysis['has_files']:
            category_name = "Document Processing"
        elif any(word in full_content for word in ['error', 'bug', 'fix', 'problem', 'trouble']):
            category_name = "Troubleshooting"
        elif any(word in full_content for word in ['report', 'business', 'professional']):
            category_name = "Business Reports"
        
        return category_name, content_analysis
        
    except Exception as e:
        logger.warning(f"Error categorizing conversation: {e}")
        return "Help & Questions", {'has_code': False, 'has_files': False, 'has_analysis': False}

def generate_chat_name(messages):
    """Generate a chat name using enhanced analysis"""
    try:
        # Get the first user message
        user_msg = ""
        for msg in messages:
            if msg.get('role') == 'user':
                user_msg = msg.get('content', '').strip()
                break
        
        if not user_msg:
            return f"Chat {time.strftime('%m/%d %H:%M')}"
        
        # Use categorization to help with naming
        category_name, analysis = categorize_conversation(messages)
        
        # Generate name based on category and content
        user_msg_lower = user_msg.lower()
        
        if category_name == "Data Analysis":
            if 'csv' in user_msg_lower or 'excel' in user_msg_lower:
                return "Excel/CSV Data Analysis"
            else:
                return "Data Analysis Project"
        elif category_name == "Business Reports":
            return "Business Report Generation"
        elif category_name == "Code Development":
            if 'python' in analysis.get('code_types', []):
                return "Python Development"
            else:
                return "Code Development"
        elif category_name == "Document Processing":
            return "Document Analysis"
        elif category_name == "Presentations":
            return "Chart & Visualization"
        elif category_name == "Troubleshooting":
            return "Problem Solving"
        else:
            # Generate from user message
            if len(user_msg) > 50:
                words = user_msg.split()[:4]
                return ' '.join(words).title()
            elif len(user_msg) <= 30:
                return user_msg.title()
            else:
                return user_msg[:27] + "..."
        
    except Exception as e:
        logger.warning(f"Error generating chat name: {e}")
        return f"Chat {time.strftime('%m/%d %H:%M')}"

# --- Helper Functions ---
ALLOWED_EXTENSIONS = {
    'txt', 'csv', 'jsonl',  # Text files
    'pdf',                   # PDF files
    'docx', 'doc',          # Word documents
    'xlsx', 'xls',          # Excel spreadsheets
    'pptx', 'ppt',          # PowerPoint presentations
    'png', 'jpg', 'jpeg'    # Images (for converted documents)
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_docx(file_bytes):
    """Extract text from Word document"""
    if not DOCX_AVAILABLE:
        logger.error("python-docx not available for Word document processing")
        return None
    
    try:
        doc = Document(io.BytesIO(file_bytes))
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from Word document: {e}")
        return None

def extract_text_from_excel(file_bytes):
    """Extract text from Excel spreadsheet"""
    if not OPENPYXL_AVAILABLE:
        logger.error("openpyxl not available for Excel document processing")
        return None
    
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        text_content = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content.append(f"=== Sheet: {sheet_name} ===")
            
            for row in sheet.iter_rows(values_only=True):
                row_data = []
                for cell in row:
                    if cell is not None:
                        row_data.append(str(cell))
                if any(cell.strip() for cell in row_data if cell):
                    text_content.append(" | ".join(row_data))
        
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from Excel document: {e}")
        return None

def extract_text_from_pptx(file_bytes):
    """Extract text from PowerPoint presentation"""
    if not PPTX_AVAILABLE:
        logger.error("python-pptx not available for PowerPoint document processing")
        return None
    
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            text_content.append(f"=== Slide {i} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
        
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint document: {e}")
        return None

def extract_text_from_pdf(file_bytes):
    """Extract text from PDF, with fallback to image conversion"""
    if not PDFPLUMBER_AVAILABLE:
        logger.error("pdfplumber not available for PDF text extraction")
        return None, None
    
    try:
        # Try to extract text first
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            text_content = []
            has_text = False
            
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    has_text = True
                    text_content.append(f"=== Page {page_num} ===")
                    text_content.append(page_text.strip())
            
            # If we found text, return it
            if has_text and text_content:
                return "\n".join(text_content), None
            
            # If no text found, convert to images
            logger.info("No text found in PDF, converting to images")
            return convert_pdf_to_images(file_bytes)
            
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        return convert_pdf_to_images(file_bytes)

def convert_pdf_to_images(file_bytes):
    """Convert PDF to images when text extraction fails"""
    if not PDF2IMAGE_AVAILABLE or not PIL_AVAILABLE:
        logger.error("pdf2image or PIL not available for PDF to image conversion")
        return None, None
    
    try:
        # Convert PDF to images
        images = convert_from_bytes(file_bytes, dpi=150, fmt='png')
        image_data_list = []
        
        for i, image in enumerate(images, 1):
            # Convert PIL image to base64
            buffer = io.BytesIO()
            image.save(buffer, format='PNG')
            image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
            
            image_data_list.append({
                'page': i,
                'type': 'image',
                'format': 'png',
                'data': image_base64
            })
        
        # Return None for text and the image data
        return None, image_data_list
        
    except Exception as e:
        logger.error(f"Error converting PDF to images: {e}")
        return None, None

def convert_csv_to_jsonl(csv_content):
    """Convert CSV content to JSONL format"""
    try:
        # Parse CSV content
        csv_reader = csv.DictReader(io.StringIO(csv_content))
        jsonl_lines = []
        
        for row in csv_reader:
            # Convert each row to a JSON object
            jsonl_lines.append(json.dumps(row))
        
        # Join with newlines to create JSONL format
        jsonl_content = '\n'.join(jsonl_lines)
        return jsonl_content
    except Exception as e:
        logger.error(f"Error converting CSV to JSONL: {e}")
        return None

def convert_txt_to_jsonl(txt_content):
    """Convert TXT content to JSONL format"""
    try:
        # Split text into lines and create JSON objects
        lines = txt_content.strip().split('\n')
        jsonl_lines = []
        
        for i, line in enumerate(lines):
            if line.strip():  # Skip empty lines
                # Create a JSON object with line number and content
                json_obj = {
                    "line_number": i + 1,
                    "content": line.strip()
                }
                jsonl_lines.append(json.dumps(json_obj))
        
        # Join with newlines to create JSONL format
        jsonl_content = '\n'.join(jsonl_lines)
        return jsonl_content
    except Exception as e:
        logger.error(f"Error converting TXT to JSONL: {e}")
        return None

def convert_to_jsonl(file_content, file_extension, file_bytes=None):
    """Convert various file types to JSONL format"""
    try:
        if file_extension == '.csv':
            return convert_csv_to_jsonl(file_content)
        elif file_extension == '.txt':
            return convert_txt_to_jsonl(file_content)
        elif file_extension == '.jsonl':
            # Already JSONL, just validate and return
            lines = file_content.strip().split('\n')
            validated_lines = []
            for line in lines:
                if line.strip():
                    try:
                        # Validate JSON
                        json.loads(line.strip())
                        validated_lines.append(line.strip())
                    except json.JSONDecodeError:
                        logger.warning(f"Invalid JSON line skipped: {line[:50]}...")
            return '\n'.join(validated_lines)
        elif file_extension in ['.docx', '.doc']:
            if file_bytes is None:
                logger.error("File bytes required for Word document processing")
                return None
            extracted_text = extract_text_from_docx(file_bytes)
            if extracted_text:
                return convert_txt_to_jsonl(extracted_text)
            else:
                logger.error("Failed to extract text from Word document")
                return None
        elif file_extension in ['.xlsx', '.xls']:
            if file_bytes is None:
                logger.error("File bytes required for Excel document processing")
                return None
            extracted_text = extract_text_from_excel(file_bytes)
            if extracted_text:
                return convert_txt_to_jsonl(extracted_text)
            else:
                logger.error("Failed to extract text from Excel document")
                return None
        elif file_extension in ['.pptx', '.ppt']:
            if file_bytes is None:
                logger.error("File bytes required for PowerPoint document processing")
                return None
            extracted_text = extract_text_from_pptx(file_bytes)
            if extracted_text:
                return convert_txt_to_jsonl(extracted_text)
            else:
                logger.error("Failed to extract text from PowerPoint document")
                return None
        elif file_extension == '.pdf':
            if file_bytes is None:
                logger.error("File bytes required for PDF document processing")
                return None
            extracted_text, image_data = extract_text_from_pdf(file_bytes)
            if extracted_text:
                # PDF had extractable text
                return convert_txt_to_jsonl(extracted_text)
            elif image_data:
                # PDF was converted to images
                jsonl_lines = []
                for img_data in image_data:
                    json_obj = {
                        "page": img_data['page'],
                        "type": "image",
                        "format": img_data['format'],
                        "data": img_data['data'],
                        "description": f"Page {img_data['page']} of PDF document (converted to image)"
                    }
                    jsonl_lines.append(json.dumps(json_obj))
                return '\n'.join(jsonl_lines)
            else:
                logger.error("Failed to process PDF document")
                return None
        else:
            # For other file types, create a simple JSONL structure
            json_obj = {
                "filename": "uploaded_file",
                "content": file_content,
                "file_type": file_extension.replace('.', '')
            }
            return json.dumps(json_obj)
    except Exception as e:
        logger.error(f"Error converting {file_extension} to JSONL: {e}")
        return None

def get_initial_system_prompt():
    return {
        "role": "system",
        "content": (
            "You are BusinessAstra, an intelligent AI business assistant with advanced document processing and analysis capabilities. "
            "You can help with business analytics, strategy, market research, and comprehensive document analysis. "
            "When users upload files, they are automatically processed and converted for optimal analysis:\n\n"
            "**Supported File Types:**\n"
            "- Word documents (.docx, .doc) → Text extraction\n"
            "- Excel spreadsheets (.xlsx, .xls) → Data extraction with sheet structure\n"
            "- PowerPoint presentations (.pptx, .ppt) → Text extraction from slides\n"
            "- PDF files → Smart processing: text extraction for digital PDFs, image conversion for scanned PDFs\n"
            "- Images (.png, .jpg, .jpeg) → Direct image analysis\n"
            "- Text files (.txt, .csv) → Direct content processing\n\n"
            "You have access to a powerful code execution sandbox that supports:\n"
            "- Python code execution (with libraries like matplotlib, pandas, numpy, openpyxl)\n"
            "- ReportLab for PDF document generation (professional PDFs, reports, charts)\n"
            "- Manim for creating mathematical animations and videos\n\n"
            "IMPORTANT: When providing code solutions, always include the code in proper markdown code blocks:\n"
            "- Use ```python for Python code\n"
            "- Use ```reportlab for ReportLab PDF generation code\n"
            "- Use ```manim for Manim code\n\n"
            "After providing code, mention that the user can click the play button (▶️) to execute it automatically. "
            "Generated files will be organized in:\n"
            "- /static/outputs/images/ for images (PNG, JPG, SVG)\n"
            "- /static/outputs/documents/ for PDFs and documents\n"
            "- /static/outputs/videos/ for MP4 and video files\n"
            "Always provide download links in your responses when code generates files."
        )
    }

def ensure_session_messages():
    if 'messages' not in session or not session['messages']:
        session['messages'] = [get_initial_system_prompt()]
        session.modified = True
        logger.info("Initialized new conversation history in session.")
    elif session['messages'][0]['role'] != 'system':
        session['messages'].insert(0, get_initial_system_prompt())
        session.modified = True
        logger.warning("System prompt was missing/misplaced. Re-inserted.")
    # Also initialize a place to store uploaded file info
    if 'uploaded_files_info' not in session:
        session['uploaded_files_info'] = {} # Store {filename: file_id}
        session.modified = True
    return session['messages']

def get_model_name():
    # Use the specific model you requested, or a fallback from .env, or a general default
    return os.getenv("TOGETHER_MODEL_NAME", "meta-llama/Llama-3.3-70B-Instruct-Turbo-Free")
    # return "meta-llama/Llama-3.3-70B-Instruct-Turbo-Free" # If this is the correct one

def get_files_context():
    """Get context from locally stored files - only for the first message after upload"""
    if 'uploaded_files_content' not in session or 'files_context_used' not in session:
        return ""
    
    # Check if we should include file context (only for first message after upload)
    if session.get('files_context_used', True):
        return ""
    
    files_context = []
    for filename, file_info in session['uploaded_files_content'].items():
        files_context.append(f"\n--- File: {filename} (Original type: {file_info['original_type']}) ---")
        files_context.append(file_info['content'][:2000])  # Limit content to avoid token limits
        if len(file_info['content']) > 2000:
            files_context.append("... (content truncated)")
        files_context.append("--- End of file ---\n")
    
    if files_context:
        # Mark that file context has been used
        session['files_context_used'] = True
        session.modified = True
        return "\n\nUploaded files content:\n" + "\n".join(files_context)
    return ""

# --- Sandbox Execution System ---
class PythonSandbox:
    def __init__(self):
        # Use system temp directory to avoid Flask file watcher issues
        import tempfile
        # Create sandbox in a completely isolated temp directory
        temp_root = tempfile.gettempdir()
        self.base_sandbox_dir = os.path.join(temp_root, "businessastra_sandbox_isolated")
        self.output_dir = os.path.join(os.getcwd(), "static", "outputs")
        
        # Create organized output directories
        self.images_dir = os.path.join(self.output_dir, "images")
        self.documents_dir = os.path.join(self.output_dir, "documents")
        self.videos_dir = os.path.join(self.output_dir, "videos")
        
        # Create directories if they don't exist
        os.makedirs(self.base_sandbox_dir, exist_ok=True)
        os.makedirs(self.images_dir, exist_ok=True)
        os.makedirs(self.documents_dir, exist_ok=True)
        os.makedirs(self.videos_dir, exist_ok=True)
        
        logger.info(f"Sandbox directory: {self.base_sandbox_dir}")
        logger.info(f"Output directory: {self.output_dir}")
        
        # Clean up all sandbox directories on startup (fresh start)
        self.cleanup_all_sandboxes()
        
        # Set up periodic cleanup
        self.setup_periodic_cleanup()
        
        # Timeout for code execution (in seconds)
        self.timeout = 30
        
    def create_sandbox_environment(self):
        """Create a temporary sandbox directory for code execution"""
        sandbox_id = str(uuid.uuid4())
        sandbox_path = os.path.join(self.base_sandbox_dir, sandbox_id)
        os.makedirs(sandbox_path, exist_ok=True)
        return sandbox_id, sandbox_path
    
    def cleanup_sandbox(self, sandbox_path):
        """Clean up the sandbox directory"""
        try:
            if os.path.exists(sandbox_path):
                shutil.rmtree(sandbox_path)
        except Exception as e:
            logger.warning(f"Failed to cleanup sandbox {sandbox_path}: {e}")
    
    def cleanup_old_sandboxes(self):
        """Clean up old sandbox directories (older than 10 minutes)"""
        try:
            if not os.path.exists(self.base_sandbox_dir):
                return
            
            current_time = time.time()
            cleaned_count = 0
            
            for item in os.listdir(self.base_sandbox_dir):
                item_path = os.path.join(self.base_sandbox_dir, item)
                if os.path.isdir(item_path):
                    # Check if directory is older than 10 minutes
                    creation_time = os.path.getctime(item_path)
                    if current_time - creation_time > 600:  # 10 minutes
                        try:
                            shutil.rmtree(item_path)
                            cleaned_count += 1
                            logger.debug(f"Cleaned up old sandbox: {item}")
                        except Exception as e:
                            logger.warning(f"Failed to cleanup old sandbox {item}: {e}")
            
            if cleaned_count > 0:
                logger.info(f"Cleaned up {cleaned_count} old sandbox directories")
                
        except Exception as e:
            logger.warning(f"Failed to cleanup old sandboxes: {e}")
    
    def cleanup_all_sandboxes(self):
        """Clean up all sandbox directories (for server restart)"""
        try:
            if os.path.exists(self.base_sandbox_dir):
                shutil.rmtree(self.base_sandbox_dir)
                os.makedirs(self.base_sandbox_dir, exist_ok=True)
                logger.info("Cleaned up all sandbox directories on restart")
        except Exception as e:
            logger.warning(f"Failed to cleanup all sandboxes: {e}")
    
    def setup_periodic_cleanup(self):
        """Set up periodic cleanup every 5 minutes"""
        def cleanup_worker():
            while True:
                time.sleep(300)  # 5 minutes
                self.cleanup_old_sandboxes()
        
        cleanup_thread = threading.Thread(target=cleanup_worker, daemon=True)
        cleanup_thread.start()
        logger.info("Started periodic sandbox cleanup (every 5 minutes)")
    
    def execute_python_code(self, code, sandbox_id=None, session_data=None):
        """Execute Python code in a sandbox environment"""
        if not sandbox_id:
            sandbox_id, sandbox_path = self.create_sandbox_environment()
        else:
            sandbox_path = os.path.join(self.base_sandbox_dir, sandbox_id)
            
        try:
            # Copy uploaded files to sandbox if they exist in session
            self._copy_uploaded_files_to_sandbox(sandbox_path, session_data)
            
            # Process the code to fix file paths
            processed_code = self._process_code_paths(code)
            
            # Create the Python file with matplotlib backend fix
            code_file = os.path.join(sandbox_path, "main.py")
            
            # Add matplotlib backend configuration to prevent popup windows
            matplotlib_fix = """
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
plt.ioff()  # Turn off interactive mode

"""
            
            # Check if the code uses matplotlib and add the fix
            if any(lib in processed_code.lower() for lib in ['matplotlib', 'plt.', 'pyplot']):
                final_code = matplotlib_fix + processed_code
            else:
                final_code = processed_code
            
            with open(code_file, 'w', encoding='utf-8') as f:
                f.write(final_code)
            
            # Execute the code with timeout and isolated environment
            # Create a clean environment to prevent interference
            clean_env = os.environ.copy()
            clean_env['PYTHONPATH'] = sandbox_path
            clean_env['PYTHONDONTWRITEBYTECODE'] = '1'  # Prevent .pyc file creation
            clean_env['PYTHONUNBUFFERED'] = '1'        # Prevent buffering issues
            
            result = subprocess.run(
                [sys.executable, code_file],
                cwd=sandbox_path,
                capture_output=True,
                text=True,
                timeout=self.timeout,
                env=clean_env
            )
            
            # Collect output files
            output_files = self._collect_output_files(sandbox_path, sandbox_id)
            
            return {
                'success': result.returncode == 0,
                'stdout': result.stdout,
                'stderr': result.stderr,
                'return_code': result.returncode,
                'output_files': output_files,
                'sandbox_id': sandbox_id
            }
            
        except subprocess.TimeoutExpired:
            return {
                'success': False,
                'stdout': '',
                'stderr': f'Code execution timed out after {self.timeout} seconds',
                'return_code': -1,
                'output_files': [],
                'sandbox_id': sandbox_id
            }
        except Exception as e:
            return {
                'success': False,
                'stdout': '',
                'stderr': f'Execution error: {str(e)}',
                'return_code': -1,
                'output_files': [],
                'sandbox_id': sandbox_id
            }
    
    def execute_reportlab_code(self, python_code, sandbox_id=None):
        """Execute Python code with ReportLab for PDF generation"""
        if not REPORTLAB_AVAILABLE:
            return {
                'success': False,
                'stdout': '',
                'stderr': 'ReportLab not available. Please install: pip install reportlab',
                'return_code': -1,
                'output_files': [],
                'sandbox_id': sandbox_id or 'unknown'
            }
            
        if not sandbox_id:
            sandbox_id, sandbox_path = self.create_sandbox_environment()
        else:
            sandbox_path = os.path.join(self.base_sandbox_dir, sandbox_id)
            
        try:
            # Copy uploaded files to sandbox if they exist in session
            self._copy_uploaded_files_to_sandbox(sandbox_path, session_data=session)
            
            # Process the code to fix file paths
            processed_code = self._process_code_paths(python_code)
            
            # Create the Python file with ReportLab imports
            code_file = os.path.join(sandbox_path, "main.py")
            
            # Add ReportLab imports and setup
            reportlab_setup = """
import os
import sys
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors

"""
            
            final_code = reportlab_setup + processed_code
            
            with open(code_file, 'w', encoding='utf-8') as f:
                f.write(final_code)
            
            # Execute the code with timeout and isolated environment
            clean_env = os.environ.copy()
            clean_env['PYTHONPATH'] = sandbox_path
            clean_env['PYTHONDONTWRITEBYTECODE'] = '1'
            clean_env['PYTHONUNBUFFERED'] = '1'
            
            result = subprocess.run(
                [sys.executable, code_file],
                cwd=sandbox_path,
                capture_output=True,
                text=True,
                timeout=self.timeout,
                env=clean_env
            )
            
            # Collect output files (especially PDFs)
            output_files = self._collect_output_files(sandbox_path, sandbox_id, file_types=['.pdf', '.png', '.jpg'])
            
            return {
                'success': result.returncode == 0,
                'stdout': result.stdout,
                'stderr': result.stderr,
                'return_code': result.returncode,
                'output_files': output_files,
                'sandbox_id': sandbox_id
            }
            
        except subprocess.TimeoutExpired:
            return {
                'success': False,
                'stdout': '',
                'stderr': f'ReportLab PDF generation timed out after {self.timeout} seconds',
                'return_code': -1,
                'output_files': [],
                'sandbox_id': sandbox_id
            }
        except Exception as e:
            return {
                'success': False,
                'stdout': '',
                'stderr': f'ReportLab PDF generation error: {str(e)}',
                'return_code': -1,
                'output_files': [],
                'sandbox_id': sandbox_id
            }
    
    def execute_manim_code(self, python_code, sandbox_id=None):
        """Execute Manim code to generate video"""
        if not sandbox_id:
            sandbox_id, sandbox_path = self.create_sandbox_environment()
        else:
            sandbox_path = os.path.join(self.base_sandbox_dir, sandbox_id)
            
        try:
            # Create the Manim Python file
            manim_file = os.path.join(sandbox_path, "manim_scene.py")
            with open(manim_file, 'w', encoding='utf-8') as f:
                f.write(python_code)
            
            # Run manim with isolated environment
            clean_env = os.environ.copy()
            clean_env['PYTHONPATH'] = sandbox_path
            clean_env['PYTHONDONTWRITEBYTECODE'] = '1'  # Prevent .pyc file creation
            clean_env['PYTHONUNBUFFERED'] = '1'        # Prevent buffering issues
            
            result = subprocess.run(
                ['manim', 'manim_scene.py', '-ql'],  # -ql for low quality, faster rendering
                cwd=sandbox_path,
                capture_output=True,
                text=True,
                timeout=60,  # Longer timeout for video rendering
                env=clean_env
            )
            
            # Collect output files
            output_files = self._collect_output_files(sandbox_path, sandbox_id, file_types=['.mp4', '.gif'])
            
            return {
                'success': result.returncode == 0,
                'stdout': result.stdout,
                'stderr': result.stderr,
                'return_code': result.returncode,
                'output_files': output_files,
                'sandbox_id': sandbox_id
            }
            
        except subprocess.TimeoutExpired:
            return {
                'success': False,
                'stdout': '',
                'stderr': 'Manim rendering timed out after 60 seconds',
                'return_code': -1,
                'output_files': [],
                'sandbox_id': sandbox_id
            }
        except Exception as e:
            return {
                'success': False,
                'stdout': '',
                'stderr': f'Manim execution error: {str(e)}',
                'return_code': -1,
                'output_files': [],
                'sandbox_id': sandbox_id
            }
    
    def _collect_output_files(self, sandbox_path, sandbox_id, file_types=None):
        """Collect output files from sandbox and move them to organized output directories"""
        if file_types is None:
            file_types = ['.png', '.jpg', '.jpeg', '.gif', '.mp4', '.pdf', '.svg', '.txt', '.csv', '.json']
        
        output_files = []
        
        # Define file type categories
        image_types = ['.png', '.jpg', '.jpeg', '.gif', '.svg']
        document_types = ['.pdf', '.txt', '.csv', '.json']
        video_types = ['.mp4', '.avi', '.mov', '.webm']
        
        try:
            for root, dirs, files in os.walk(sandbox_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    file_ext = os.path.splitext(file)[1].lower()
                    
                    if file_ext in file_types and file != "main.py" and file != "manim_scene.py" and file != "document.tex":
                        # Determine target directory based on file type
                        if file_ext in image_types:
                            target_dir = self.images_dir
                            url_path = f'/static/outputs/images/{sandbox_id}_{file}'
                        elif file_ext in document_types:
                            target_dir = self.documents_dir
                            url_path = f'/static/outputs/documents/{sandbox_id}_{file}'
                        elif file_ext in video_types:
                            target_dir = self.videos_dir
                            url_path = f'/static/outputs/videos/{sandbox_id}_{file}'
                        else:
                            target_dir = self.output_dir
                            url_path = f'/static/outputs/{sandbox_id}_{file}'
                        
                        # Generate unique filename
                        output_filename = f"{sandbox_id}_{file}"
                        output_path = os.path.join(target_dir, output_filename)
                        
                        # Copy file to appropriate directory
                        shutil.copy2(file_path, output_path)
                        
                        output_files.append({
                            'filename': output_filename,
                            'original_name': file,
                            'type': file_ext[1:],  # Remove the dot
                            'path': url_path,
                            'size': os.path.getsize(output_path),
                            'category': 'image' if file_ext in image_types else 'document' if file_ext in document_types else 'video' if file_ext in video_types else 'other'
                        })
        except Exception as e:
            logger.warning(f"Error collecting output files: {e}")
        
        return output_files
    
    def _copy_uploaded_files_to_sandbox(self, sandbox_path, session_data=None):
        """Copy uploaded files from session to sandbox for code execution"""
        try:
            if session_data and 'uploaded_files_content' in session_data:
                for filename, file_info in session_data['uploaded_files_content'].items():
                    # Write the file content to sandbox
                    file_path = os.path.join(sandbox_path, filename)
                    with open(file_path, 'w', encoding='utf-8') as f:
                        f.write(file_info['content'])
                    logger.info(f"Copied uploaded file {filename} to sandbox")
        except Exception as e:
            logger.warning(f"Error copying uploaded files to sandbox: {e}")
    
    def _process_code_paths(self, code):
        """Process code to fix file paths for sandbox execution"""
        import re
        
        # Replace absolute paths to /static/outputs/ with relative paths
        # This allows the code to save files in the sandbox directory
        processed_code = re.sub(
            r"['\"]\/static\/outputs\/[^'\"]*\/([^'\"]+)['\"]",
            r"'\1'",
            code
        )
        
        # Also handle cases where people use /static/outputs/ directly
        processed_code = re.sub(
            r"['\"]\/static\/outputs\/([^'\"]+)['\"]",
            r"'\1'",
            processed_code
        )
        
        return processed_code

# Sandbox will be initialized in main block

# Initialize database
if SQLALCHEMY_AVAILABLE and db:
    with app.app_context():
        init_database()

# --- Error Handlers and Health Checks ---
@app.errorhandler(404)
def not_found_error(error):
    """Handle 404 errors"""
    logger.warning(f"404 error: {request.url}")
    if request.accept_mimetypes.accept_json and not request.accept_mimetypes.accept_html:
        return jsonify({'error': 'Not found'}), 404
    return render_template('404.html'), 404

@app.errorhandler(500)
def internal_error(error):
    """Handle 500 errors"""
    logger.error(f"Internal error: {str(error)}")
    if db:
        db.session.rollback()
    if request.accept_mimetypes.accept_json and not request.accept_mimetypes.accept_html:
        return jsonify({'error': 'Internal server error'}), 500
    return render_template('500.html'), 500

@app.errorhandler(429)
def ratelimit_handler(e):
    """Handle rate limit errors"""
    logger.warning(f"Rate limit exceeded: {request.remote_addr}")
    return jsonify({
        'error': 'Rate limit exceeded',
        'message': str(e.description),
        'retry_after': getattr(e, 'retry_after', None)
    }), 429

@app.route('/health')
def health_check():
    """Health check endpoint for monitoring"""
    try:
        # Check database connectivity
        db_status = 'ok'
        if SQLALCHEMY_AVAILABLE and db:
            try:
                db.session.execute(text('SELECT 1'))
                db.session.commit()
            except Exception as e:
                db_status = f'error: {str(e)}'
                logger.error(f"Database health check failed: {e}")
        
        # Check API connectivity
        api_status = 'ok' if TOGETHER_CHAT_CLIENT else 'not configured'
        
        health_data = {
            'status': 'healthy' if db_status == 'ok' else 'unhealthy',
            'timestamp': datetime.datetime.utcnow().isoformat(),
            'version': '1.0.0',
            'components': {
                'database': db_status,
                'together_api': api_status,
                'document_processing': {
                    'docx': DOCX_AVAILABLE,
                    'excel': OPENPYXL_AVAILABLE,
                    'pptx': PPTX_AVAILABLE,
                    'pdf': PDFPLUMBER_AVAILABLE,
                    'reportlab': REPORTLAB_AVAILABLE
                }
            }
        }
        
        status_code = 200 if health_data['status'] == 'healthy' else 503
        return jsonify(health_data), status_code
        
    except Exception as e:
        logger.error(f"Health check failed: {e}")
        return jsonify({
            'status': 'unhealthy',
            'error': str(e),
            'timestamp': datetime.datetime.utcnow().isoformat()
        }), 503

@app.route('/metrics')
def metrics():
    """Basic metrics endpoint"""
    try:
        metrics_data = {
            'timestamp': datetime.datetime.utcnow().isoformat(),
            'uptime': time.time() - app.config.get('_start_time', time.time()),
        }
        
        if SQLALCHEMY_AVAILABLE and db:
            try:
                # Get basic database stats
                with db.engine.connect() as conn:
                    if hasattr(db, 'Model'):  # Check if models are available
                        user_count = User.query.count() if 'User' in globals() else 0
                        conversation_count = Conversation.query.count() if 'Conversation' in globals() else 0
                        metrics_data['database'] = {
                            'users': user_count,
                            'conversations': conversation_count
                        }
            except Exception as e:
                logger.warning(f"Could not get database metrics: {e}")
        
        return jsonify(metrics_data)
        
    except Exception as e:
        logger.error(f"Metrics endpoint failed: {e}")
        return jsonify({'error': 'Metrics unavailable'}), 500

# Store start time for uptime calculation
app.config['_start_time'] = time.time()

# --- Flask Routes ---

# Authentication Routes
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')
        
        if not email or not password:
            flash('Please provide both email and password', 'error')
            return render_template('login.html')
        
        if verify_user(email, password):
            session['user_email'] = email
            flash('Login successful!', 'success')
            return redirect(url_for('index'))
        else:
            flash('Invalid email or password', 'error')
            return render_template('login.html')
    
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form.get('username')
        email = request.form.get('email')
        password = request.form.get('password')
        confirm_password = request.form.get('confirm_password')
        
        if not username or not email or not password or not confirm_password:
            flash('Please fill in all fields', 'error')
            return render_template('register.html')
        
        if len(username) < 3:
            flash('Username must be at least 3 characters long', 'error')
            return render_template('register.html')
        
        if password != confirm_password:
            flash('Passwords do not match', 'error')
            return render_template('register.html')
        
        if len(password) < 6:
            flash('Password must be at least 6 characters long', 'error')
            return render_template('register.html')
        
        success, message = create_user(email, password, username)
        if success:
            flash('Registration successful! Please log in.', 'success')
            return redirect(url_for('login'))
        else:
            flash(message, 'error')
            return render_template('register.html')
    
    return render_template('register.html')

@app.route('/logout')
def logout():
    session.pop('user_email', None)
    flash('You have been logged out', 'info')
    return redirect(url_for('login'))

@app.route('/')
@login_required
def index():
    ensure_session_messages()
    
    # Get user information
    users = load_users()
    user_email = session.get('user_email')
    user_data = users.get(user_email, {})
    username = user_data.get('username', 'User')
    
    # Get chat history
    chat_history = get_user_chat_history(user_email)
    
    # Get current session messages (excluding system prompt for display)
    current_messages = session.get('messages', [])
    display_messages = [msg for msg in current_messages if msg.get('role') != 'system']
    
    # Get current chat info
    current_chat_id = session.get('current_chat_id')
    current_chat_name = None
    if current_chat_id:
        for chat in chat_history:
            if chat.get('id') == current_chat_id:
                current_chat_name = chat.get('name', 'Unnamed Chat')
                break
    
    return render_template('index.html', 
                         username=username, 
                         chat_history=chat_history,
                         current_messages=display_messages,
                         current_chat_name=current_chat_name)

@app.route('/execute_code', methods=['POST'])
@login_required
def execute_code():
    """Execute Python code in sandbox"""
    try:
        data = request.json
        code = data.get('code', '')
        execution_type = data.get('type', 'python')  # python, reportlab, manim
        
        if not code.strip():
            return jsonify({"error": "No code provided"}), 400
        
        logger.info(f"Executing {execution_type} code in sandbox")
        
        if execution_type == 'python':
            result = sandbox.execute_python_code(code, session_data=session)
        elif execution_type == 'reportlab':
            result = sandbox.execute_reportlab_code(code)
        elif execution_type == 'manim':
            result = sandbox.execute_manim_code(code)
        else:
            return jsonify({"error": f"Unsupported execution type: {execution_type}"}), 400
        
        # Don't add to chat automatically - frontend will handle this
        
        return jsonify({
            "success": result['success'],
            "stdout": result['stdout'],
            "stderr": result['stderr'],
            "output_files": result['output_files'],
            "message": "Code executed successfully" if result['success'] else "Code execution failed"
        }), 200
        
    except Exception as e:
        logger.error(f"Error in code execution: {e}")
        return jsonify({"error": f"Execution error: {str(e)}"}), 500

@app.route('/static/outputs/<path:filename>')
@login_required
def serve_output_file(filename):
    """Serve generated output files from organized directories"""
    try:
        output_dir = os.path.join(os.getcwd(), "static", "outputs")
        file_path = os.path.join(output_dir, filename)
        
        if os.path.exists(file_path):
            return send_file(file_path)
        else:
            return jsonify({"error": "File not found"}), 404
    except Exception as e:
        logger.error(f"Error serving file {filename}: {e}")
        return jsonify({"error": "File not found"}), 404

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    # Enhanced rate limiting for chat endpoint
    if limiter:
        try:
            limiter.check()
        except Exception as e:
            return jsonify({
                'error': 'Rate limit exceeded for chat',
                'message': 'Please wait before sending another message'
            }), 429
    if not TOGETHER_CHAT_CLIENT:
        logger.error("Chat attempt: Together AI CHAT client unavailable.")
        return jsonify({"error": "AI service is not configured."}), 503

    # Check rate limiting
    estimated_tokens = 1500  # Conservative estimate
    can_proceed, rate_message = rate_limiter.can_make_request(estimated_tokens)
    
    if not can_proceed:
        wait_time = rate_limiter.get_wait_time()
        logger.warning(f"Rate limit exceeded: {rate_message}")
        return jsonify({
            "error": f"Rate limit exceeded. {rate_message}. Please wait {int(wait_time)} seconds.",
            "wait_time": int(wait_time)
        }), 429

    current_messages = ensure_session_messages()
    data = request.json
    user_input = data.get('message')

    if not user_input:
        return jsonify({"error": "No message provided"}), 400

    # Add file context to user message if files are available
    files_context = get_files_context()
    user_message_with_context = user_input + files_context
    
    current_messages.append({"role": "user", "content": user_message_with_context})
    session.modified = True
    logger.info(f"User: '{user_input[:70]}...'")

    # Check if the user is referring to an uploaded file by name
    # This is a very basic check. More sophisticated NLP could be used.
    for filename, file_id in session.get('uploaded_files_info', {}).items():
        if filename.split('.')[0].lower() in user_input.lower(): # Simple check if filename (sans ext) is in input
            context_note = f"(System note: User might be referring to the uploaded file '{filename}' with ID '{file_id}'. This file was uploaded to Together AI.)"
            if not any(msg['role'] == 'system' and file_id in msg['content'] for msg in current_messages[-3:]): # Avoid spamming
                 current_messages.insert(-1, {"role": "system", "content": context_note}) # Insert before last user message
                 session.modified = True
                 logger.info(f"Added context note for file '{filename}' (ID: {file_id}) based on user input.")
            break
    logger.debug(f"History before API: {current_messages}")


    def generate_response():
        try:
            # Record the API request
            rate_limiter.record_request(estimated_tokens)
            
            full_response_content = ""
            model_to_use = get_model_name()
            stream = TOGETHER_CHAT_CLIENT.chat.completions.create(
                model=model_to_use,
                messages=current_messages,
                stream=True,
                max_tokens=1500,
            )
            for chunk in stream:
                if hasattr(chunk, "choices") and chunk.choices[0].delta and chunk.choices[0].delta.content:
                    content = chunk.choices[0].delta.content
                    yield content
                    full_response_content += content
            
            if full_response_content:
                current_messages.append({"role": "assistant", "content": full_response_content})
                session.modified = True
                logger.info(f"AI: '{full_response_content[:70]}...'")
        except Exception as e:
            error_str = str(e).lower()
            logger.error(f"API call error with model {model_to_use}: {e}", exc_info=True)
            
            # Handle specific rate limit errors
            if "rate limit" in error_str or "too many requests" in error_str:
                yield "⚠️ Rate limit reached. Please wait a moment before sending another message."
            elif "timeout" in error_str:
                yield "⚠️ Request timed out. Please try again."
            elif "connection" in error_str:
                yield "⚠️ Connection error. Please check your internet and try again."
            else:
                yield f"⚠️ AI service temporarily unavailable. Please try again in a moment."
            
            # Remove the user message on error
            if current_messages and current_messages[-1]['role'] == 'user':
                current_messages.pop()
                session.modified = True
    return Response(stream_with_context(generate_response()), mimetype='text/plain; charset=utf-8')

@app.route('/upload', methods=['POST'])
@login_required
def upload_file_to_together():
    # Rate limiting for file uploads
    if limiter:
        try:
            limiter.check()
        except Exception as e:
            return jsonify({
                'error': 'Upload rate limit exceeded',
                'message': 'Please wait before uploading another file'
            }), 429
    if not TOGETHER_API_KEY:
        logger.error("Upload attempt: Together API key not available.")
        return jsonify({"error": "AI service is not configured."}), 503
    
    # Debug API key (show only first and last few characters for security)
    api_key_preview = f"{TOGETHER_API_KEY[:8]}...{TOGETHER_API_KEY[-8:]}" if len(TOGETHER_API_KEY) > 16 else "***"
    logger.debug(f"Using API key: {api_key_preview}")
        
    current_messages = ensure_session_messages() # To add system notes

    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    flask_file = request.files['file']
    if flask_file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if flask_file and allowed_file(flask_file.filename):
        original_filename = secure_filename(flask_file.filename)
        
        # Temporarily save the file to pass its path or read its content for the `requests` library
        # Or, read it into memory directly if `requests` handles in-memory file objects well for multipart
        
        file_bytes = flask_file.read() # Read the file content into bytes
        flask_file.seek(0) # Reset pointer if you need to read it again or save it

        # Convert all files to JSONL format for better compatibility with Gemini
        file_extension = '.' + original_filename.rsplit('.', 1)[1].lower()
        try:
            # Handle different file types
            file_content = None
            
            if file_extension in ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.pdf']:
                # Binary documents - process directly with file_bytes
                logger.info(f"Processing {file_extension} document: {original_filename}")
                jsonl_content = convert_to_jsonl(None, file_extension, file_bytes)
            elif file_extension in ['.png', '.jpg', '.jpeg']:
                # Image files - convert to base64
                logger.info(f"Processing image file: {original_filename}")
                image_base64 = base64.b64encode(file_bytes).decode('utf-8')
                json_obj = {
                    "type": "image",
                    "format": file_extension.replace('.', ''),
                    "data": image_base64,
                    "filename": original_filename
                }
                jsonl_content = json.dumps(json_obj)
            else:
                # Text-based files - decode as UTF-8 first
                try:
                    file_content = file_bytes.decode('utf-8')
                    jsonl_content = convert_to_jsonl(file_content, file_extension)
                except UnicodeDecodeError:
                    logger.error(f"Could not decode {file_extension} file as UTF-8")
                    return jsonify({"error": f"File {original_filename} is not a valid text file"}), 400
            
            if jsonl_content is None:
                return jsonify({"error": f"Failed to convert {file_extension} to JSONL format"}), 400
            
            # Update file data for JSONL
            processed_file_bytes = jsonl_content.encode('utf-8')
            
            # Change filename to indicate conversion (unless already JSONL)
            if not original_filename.endswith('.jsonl'):
                base_name = original_filename.rsplit('.', 1)[0]
                if file_extension in ['.png', '.jpg', '.jpeg']:
                    original_filename = f"{base_name}_image.jsonl"
                else:
                    original_filename = f"{base_name}_converted.jsonl"
                logger.info(f"Converted {file_extension} to JSONL format. New filename: {original_filename}")
            
        except Exception as e:
            logger.error(f"Error processing {file_extension} file: {e}")
            return jsonify({"error": f"Failed to process {file_extension} file: {str(e)}"}), 400

        # File processing is now done locally - no need for Together AI upload preparation

        # Process file locally instead of uploading to Together AI
        logger.info(f"Processing file '{original_filename}' locally. Original type: {file_extension}")
        
        # Store file content locally for chat context
        if 'uploaded_files_content' not in session:
            session['uploaded_files_content'] = {}
        
        # Store the converted JSONL content for use in chat
        session['uploaded_files_content'][original_filename] = {
            'content': jsonl_content,
            'original_type': file_extension,
            'upload_time': str(int(time.time())),
            'size': len(processed_file_bytes)
        }
        
        # Reset the context flag so file content will be included in the next message
        session['files_context_used'] = False
        
        # Add a system message to the chat about the uploaded file
        current_messages = ensure_session_messages()
        
        # Create appropriate system message based on file type
        if file_extension == '.pdf':
            # Check if it was converted to images or text
            if 'type": "image"' in jsonl_content:
                system_message_content = f"User has uploaded a PDF file: '{original_filename}' that was converted to images (scanned PDF). The images are now available for analysis."
            else:
                system_message_content = f"User has uploaded a PDF file: '{original_filename}' and text was extracted. The content is now available for analysis."
        elif file_extension in ['.docx', '.doc']:
            system_message_content = f"User has uploaded a Word document: '{original_filename}' and text was extracted. The content is now available for analysis."
        elif file_extension in ['.xlsx', '.xls']:
            system_message_content = f"User has uploaded an Excel spreadsheet: '{original_filename}' and data was extracted. The content is now available for analysis."
        elif file_extension in ['.pptx', '.ppt']:
            system_message_content = f"User has uploaded a PowerPoint presentation: '{original_filename}' and text was extracted. The content is now available for analysis."
        elif file_extension in ['.png', '.jpg', '.jpeg']:
            system_message_content = f"User has uploaded an image: '{original_filename}'. The image is now available for analysis."
        else:
            system_message_content = f"User has uploaded a file: '{original_filename}' (Original format: {file_extension}). The content is now available for analysis."
        
        # Add preview of content for text-based conversions
        if file_extension not in ['.png', '.jpg', '.jpeg'] and 'type": "image"' not in jsonl_content:
            preview_content = jsonl_content[:1000] + ('...' if len(jsonl_content) > 1000 else '')
            system_message_content += f"\n\nPreview:\n{preview_content}"
        
        current_messages.append({"role": "system", "content": system_message_content})
        session.modified = True
        
        logger.info(f"File '{original_filename}' processed successfully. Size: {len(processed_file_bytes)} bytes")
        
        return jsonify({
            "message": f"File '{original_filename}' uploaded and processed successfully!",
            "status": "success",
            "details": f"File has been converted from {file_extension} to JSONL format and is ready for analysis in chat.",
            "file_size": len(processed_file_bytes),
            "converted_format": "JSONL",
            "original_format": file_extension.replace('.', '').upper()
        }), 200
            
    else:
        return jsonify({"error": "File type not allowed or no file."}), 400

@app.route('/test_api_key', methods=['GET'])
@login_required
def test_api_key():
    """Test if the API key works by trying to list files"""
    if not TOGETHER_API_KEY:
        return jsonify({"error": "API key not configured"}), 503
    
    try:
        # Test with a simple API call - list files
        test_url = "https://api.together.xyz/v1/files"
        headers = {
            "accept": "application/json",
            "authorization": f"Bearer {TOGETHER_API_KEY.strip()}"
        }
        
        logger.info("Testing API key with file list endpoint...")
        response = requests.get(test_url, headers=headers)
        
        logger.debug(f"Test response status: {response.status_code}")
        logger.debug(f"Test response: {response.text}")
        
        if response.status_code == 200:
            return jsonify({
                "message": "API key is valid and working",
                "status": "success",
                "files_count": len(response.json().get('data', []))
            }), 200
        elif response.status_code == 401:
            return jsonify({
                "error": "API key is invalid or expired",
                "status": "unauthorized",
                "details": response.text
            }), 401
        else:
            return jsonify({
                "error": f"API test failed with status {response.status_code}",
                "details": response.text
            }), response.status_code
            
    except Exception as e:
        logger.error(f"Error testing API key: {e}")
        return jsonify({"error": f"API test failed: {str(e)}"}), 500

@app.route('/reset_chat', methods=['POST'])
@login_required
def reset_chat():
    try:
        # Save current chat to history if it has meaningful content
        current_messages = session.get('messages', [])
        user_email = session.get('user_email')
        
        # Check if we have a meaningful conversation (more than system prompt + 1 exchange)
        user_messages = [msg for msg in current_messages if msg.get('role') == 'user']
        
        if len(user_messages) >= 1 and user_email:
            try:
                # Generate chat name
                chat_name = generate_chat_name(current_messages)
                
                # Create unique chat ID
                chat_id = str(uuid.uuid4())
                
                # Save to history
                chat_data = {
                    'id': chat_id,
                    'name': chat_name,
                    'messages': current_messages.copy(),  # Make a copy
                    'created_at': time.time(),
                    'updated_at': time.time()
                }
                
                success = save_chat_to_history(user_email, chat_data)
                if success:
                    logger.info(f"Saved chat to history: {chat_name} (ID: {chat_id})")
                else:
                    logger.warning("Failed to save chat to history")
                    
            except Exception as e:
                logger.error(f"Error saving chat to history: {e}")
        
        # Reset session
        session['messages'] = [get_initial_system_prompt()]
        session['uploaded_files_info'] = {}
        session['current_chat_id'] = None  # Clear current chat ID
        session.modified = True
        
        logger.info("Chat session reset successfully")
        return jsonify({"status": "success", "message": "New chat started"}), 200
        
    except Exception as e:
        logger.error(f"Error in reset_chat: {e}")
        return jsonify({"status": "error", "message": "Failed to start new chat"}), 500

@app.route('/load_chat/<chat_id>', methods=['POST'])
@login_required
def load_chat(chat_id):
    """Load a specific chat from history"""
    try:
        user_email = session.get('user_email')
        if not user_email:
            return jsonify({"status": "error", "message": "User not authenticated"}), 401
        
        chat_history = get_user_chat_history(user_email)
        
        # Find the chat
        target_chat = None
        for chat in chat_history:
            if chat.get('id') == chat_id:
                target_chat = chat
                break
        
        if not target_chat:
            logger.warning(f"Chat not found: {chat_id}")
            return jsonify({"status": "error", "message": "Chat not found"}), 404
        
        # Load the chat messages
        messages = target_chat.get('messages', [])
        if not messages:
            messages = [get_initial_system_prompt()]
        
        # Update session
        session['messages'] = messages
        session['uploaded_files_info'] = {}  # Reset uploaded files for safety
        session['current_chat_id'] = chat_id
        session.modified = True
        
        logger.info(f"Loaded chat: {target_chat.get('name', 'Unnamed')} (ID: {chat_id})")
        return jsonify({
            "status": "success", 
            "message": "Chat loaded successfully",
            "chat_name": target_chat.get('name', 'Unnamed Chat')
        }), 200
        
    except Exception as e:
        logger.error(f"Error loading chat {chat_id}: {e}")
        return jsonify({"status": "error", "message": "Failed to load chat"}), 500

@app.route('/delete_chat/<chat_id>', methods=['DELETE'])
@login_required
def delete_chat(chat_id):
    """Delete a specific chat from history"""
    user_email = session.get('user_email')
    users = load_users()
    
    if user_email not in users:
        return jsonify({"status": "error", "message": "User not found"}), 404
    
    chat_history = users[user_email].get('chat_history', [])
    
    # Remove the chat
    users[user_email]['chat_history'] = [
        chat for chat in chat_history if chat.get('id') != chat_id
    ]
    
    save_users(users)
    return jsonify({"status": "success", "message": "Chat deleted successfully"}), 200

@app.route('/rename_chat/<chat_id>', methods=['POST'])
@login_required
def rename_chat(chat_id):
    """Rename a specific chat"""
    user_email = session.get('user_email')
    new_name = request.json.get('name', '').strip()
    
    if not new_name:
        return jsonify({"status": "error", "message": "Name cannot be empty"}), 400
    
    users = load_users()
    if user_email not in users:
        return jsonify({"status": "error", "message": "User not found"}), 404
    
    chat_history = users[user_email].get('chat_history', [])
    
    # Find and update the chat
    for chat in chat_history:
        if chat.get('id') == chat_id:
            chat['name'] = new_name
            chat['updated_at'] = time.time()
            save_users(users)
            return jsonify({"status": "success", "message": "Chat renamed successfully"}), 200
    
    return jsonify({"status": "error", "message": "Chat not found"}), 404

if __name__ == '__main__':
    global sandbox
    # Initialize sandbox system
    sandbox = PythonSandbox()
    
    logger.info(f"Flask app starting. Model: {get_model_name()}")
    
    # Check for production mode environment variable
    PRODUCTION_MODE = os.environ.get('FLASK_PRODUCTION', 'false').lower() == 'true'
    
    if PRODUCTION_MODE:
        logger.info("Running in PRODUCTION mode (no auto-reload)")
        app.run(debug=False, port=5001, host='0.0.0.0')
        exit()
    
    # Disable auto-reloader for sandbox files by running without debug mode
    # or with custom configuration
    import os
    
    # Option 1: Run without debug mode to prevent auto-restart
    # app.run(debug=False, port=5001)
    
    # Option 2: Run with debug but try to minimize restarts
    try:
        # Try to patch the file watcher with more comprehensive patterns
        from werkzeug._reloader import StatReloaderLoop
        
        original_trigger_reload = StatReloaderLoop.trigger_reload
        
        def patched_trigger_reload(self, filename):
            filename_str = str(filename).lower()
            
            # Don't trigger reload for these patterns
            ignore_patterns = [
                'sandbox_executions', 'static/outputs', 'uploads_temp', '__pycache__',
                'businessastra_sandbox',  # System temp sandbox directory
                'site-packages',          # Python packages
                'core.py',               # Core Python files
                '.pyc',                  # Compiled Python files
                '.pyo',                  # Optimized Python files
                '.pyd',                  # Python extension modules
                'temp',                  # Any temp directories
                'tmp',                   # Any tmp directories
                '.git',                  # Git files
                '.vscode',               # VS Code files
                'node_modules',          # Node modules
                '.log',                  # Log files
                '.cache',                # Cache files
                'uploads/',              # Upload directories
                'static/outputs/',       # Output directories
            ]
            
            # Check if filename contains any ignore pattern
            if any(pattern in filename_str for pattern in ignore_patterns):
                logger.debug(f"Ignoring file change in: {filename}")
                return
                
            # Also ignore files in system temp directory
            import tempfile
            temp_dir = tempfile.gettempdir().lower()
            if temp_dir in filename_str:
                logger.debug(f"Ignoring temp file change: {filename}")
                return
                
            return original_trigger_reload(self, filename)
        
        StatReloaderLoop.trigger_reload = patched_trigger_reload
        logger.info("Successfully patched file watcher with comprehensive ignore patterns")
        
    except Exception as e:
        logger.warning(f"Could not patch file watcher: {e}")
        logger.info("Falling back to non-debug mode to prevent restarts")
        # If patching fails, run without debug to prevent restarts
        app.run(debug=False, port=5001)
        exit()
    
    # Run the app with debug mode (file watcher is patched)
    app.run(debug=True, port=5001, use_reloader=True, reloader_type='stat')

    